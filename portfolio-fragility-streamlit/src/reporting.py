"""Cloud-compatible report generation for Portfolio Fragility Lab."""

from __future__ import annotations

from datetime import datetime
from io import BytesIO

import numpy as np
import pandas as pd


DISCLAIMER = (
    "Results are model-based estimates derived from historical data and stated assumptions. "
    "They are not forecasts, guarantees, regulatory risk measures, or investment advice."
)


def build_excel_report(result, config) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        pd.DataFrame([config.__dict__]).to_excel(writer, sheet_name="Configuration", index=False)
        result.prices.to_excel(writer, sheet_name="Prices")
        result.asset_returns.to_excel(writer, sheet_name="Returns")
        pd.Series(result.metrics, name="Value").to_excel(writer, sheet_name="Risk Metrics")
        result.risk_contributions.to_excel(writer, sheet_name="Risk Contributions", index=False)
        result.liquidity.to_excel(writer, sheet_name="Liquidity", index=False)
        result.stress_tests.to_excel(writer, sheet_name="Stress Tests", index=False)
        result.regime_summary.to_excel(writer, sheet_name="Regimes", index=False)
        result.fragility_components.to_excel(writer, sheet_name="Fragility Score", index=False)
        result.allocation_comparison.to_excel(writer, sheet_name="Allocation Comparison", index=False)
        pd.Series(result.terminal_values, name="Terminal Value").to_excel(writer, sheet_name="Terminal Values", index=False)
    return output.getvalue()


def risk_metrics_csv(result) -> bytes:
    return pd.DataFrame(
        {"Metric": list(result.metrics), "Value": list(result.metrics.values())}
    ).to_csv(index=False).encode("utf-8")


def _status(score: float) -> str:
    if score <= 25:
        return "Resilient"
    if score <= 45:
        return "Stable"
    if score <= 65:
        return "Watch"
    if score <= 80:
        return "Elevated"
    return "Critical"


def build_pdf_report(result, config) -> bytes:
    """Generate a concise institutional-style executive PDF."""
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    output = BytesIO()
    document = SimpleDocTemplate(output, pagesize=A4, rightMargin=16 * mm, leftMargin=16 * mm, topMargin=14 * mm, bottomMargin=14 * mm)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="ReportTitle", parent=styles["Title"], fontSize=22, leading=26, textColor=colors.HexColor("#111419"), alignment=TA_CENTER))
    styles.add(ParagraphStyle(name="Section", parent=styles["Heading2"], fontSize=14, textColor=colors.HexColor("#0B6E67"), spaceBefore=10, spaceAfter=7))
    styles.add(ParagraphStyle(name="Small", parent=styles["BodyText"], fontSize=8.5, leading=11, textColor=colors.HexColor("#4B5563")))
    metrics = result.metrics
    worst = result.stress_tests.loc[result.stress_tests["Net Portfolio Shock"].idxmin()]
    driver = result.fragility_components.loc[result.fragility_components["Weighted Contribution"].idxmax(), "Component"]
    story = [
        Paragraph("Portfolio Fragility Lab", styles["ReportTitle"]),
        Paragraph("Institutional Portfolio Risk Summary", styles["Heading3"]),
        Paragraph(f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} | {', '.join(config.tickers)} | Benchmark: {config.benchmark}", styles["Small"]),
        Spacer(1, 8),
        Paragraph("Executive summary", styles["Section"]),
    ]
    kpis = [
        ["Fragility Score", "Expected Terminal", "Loss Probability", f"{config.confidence:.0%} VaR"],
        [f"{metrics['fragility_score']:.0f}/100 ({_status(metrics['fragility_score'])})", f"${metrics['expected_terminal']:,.0f}", f"{metrics['probability_loss']:.1%}", f"${metrics['var_currency']:,.0f}"],
    ]
    table = Table(kpis, colWidths=[42 * mm] * 4)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#111419")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("BACKGROUND", (0, 1), (-1, 1), colors.HexColor("#F3F4F6")), ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#A8B0BC")), ("FONTSIZE", (0, 0), (-1, -1), 8.5),
        ("TOPPADDING", (0, 0), (-1, -1), 7), ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
    ]))
    story.extend([
        table,
        Spacer(1, 8),
        Paragraph(
            f"The modeled portfolio is classified as <b>{_status(metrics['fragility_score'])}</b>. "
            f"Its main fragility driver is <b>{driver}</b>. Crisis correlation rises from "
            f"{metrics['normal_correlation']:.2f} to {metrics['crisis_correlation']:.2f}, reducing modeled diversification benefits.",
            styles["BodyText"],
        ),
        Paragraph("Allocation", styles["Section"]),
    ])
    allocation_rows = [["Ticker", "Current", "Crisis-Resilient", "Change"]]
    for ticker, current, resilient in zip(config.tickers, config.weights, result.resilient_weights):
        allocation_rows.append([ticker, f"{current:.1%}", f"{resilient:.1%}", f"{resilient-current:+.1%}"])
    allocation_table = Table(allocation_rows, colWidths=[42 * mm] * 4)
    allocation_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#252B33")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#A8B0BC")), ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("FONTSIZE", (0, 0), (-1, -1), 8.5), ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F6F7")]),
    ]))
    story.extend([
        allocation_table,
        Paragraph("Main stress scenario", styles["Section"]),
        Paragraph(
            f"<b>{worst['Scenario']}</b> produces a modeled net shock of {worst['Net Portfolio Shock']:.1%}, "
            f"including an estimated liquidation cost of {worst['Liquidation Cost']:.2%}. "
            f"The stressed portfolio value is ${worst['Stressed Value']:,.0f}.", styles["BodyText"]
        ),
        Paragraph("Resilient allocation comparison", styles["Section"]),
    ])
    comparison = result.allocation_comparison.copy()
    comparison_rows = [["Portfolio", "Return", "Volatility", "Worst Stress", "Liquidity"]]
    for _, row in comparison.iterrows():
        comparison_rows.append([
            row["Portfolio"], f"{row['Historical Return']:.1%}", f"{row['Historical Volatility']:.1%}",
            f"{row['Worst Stress Loss']:.1%}", f"{row['Weighted Liquidity Score']:.1f}",
        ])
    comparison_table = Table(comparison_rows, colWidths=[34 * mm] * 5)
    comparison_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#252B33")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#A8B0BC")), ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("FONTSIZE", (0, 0), (-1, -1), 8.2),
    ]))
    story.extend([
        comparison_table,
        PageBreak(),
        Paragraph("Methodology and assumptions", styles["Section"]),
        Paragraph(
            f"Method: {config.method}; paths: {config.simulations:,}; horizon: {config.simulation_days} trading days; "
            f"confidence: {config.confidence:.0%}; historical start: {config.start_date}; risk-free rate: {config.risk_free_rate:.2%}. "
            f"Liquidity events use a {config.redemption_probability:.0%} conditional redemption probability, "
            f"{config.redemption_pct:.0%} redemption size, {config.margin_call_pct:.0%} synthetic collateral call, "
            f"and {config.cash_buffer_pct:.0%} cash buffer.", styles["BodyText"]
        ),
        Paragraph("Limitations", styles["Section"]),
        Paragraph(DISCLAIMER, styles["BodyText"]),
    ])
    document.build(story)
    return output.getvalue()


def build_powerpoint_report(result, config) -> bytes:
    """Generate a compact executive PowerPoint summary."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    bg = RGBColor(5, 5, 5)
    white = RGBColor(243, 244, 246)
    muted = RGBColor(168, 176, 188)
    teal = RGBColor(34, 199, 184)

    def add_slide(title: str, lines: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = title
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        title_run.font.color.rgb = white
        body = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(11.7), Inches(5.2)).text_frame
        body.clear()
        for index, line in enumerate(lines):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = teal if index == 0 else muted
            paragraph.space_after = Pt(14)
        return slide

    metrics = result.metrics
    driver = result.fragility_components.loc[result.fragility_components["Weighted Contribution"].idxmax(), "Component"]
    worst = result.stress_tests.loc[result.stress_tests["Net Portfolio Shock"].idxmin()]
    add_slide("Portfolio Fragility Lab", [
        f"{', '.join(config.tickers)} | Benchmark: {config.benchmark}",
        f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        DISCLAIMER,
    ])
    add_slide("Executive risk summary", [
        f"Fragility Score: {metrics['fragility_score']:.0f}/100 — {_status(metrics['fragility_score'])}",
        f"Expected terminal wealth: ${metrics['expected_terminal']:,.0f}",
        f"Loss probability: {metrics['probability_loss']:.1%} | Target probability: {metrics['probability_target']:.1%}",
        f"{config.confidence:.0%} VaR: ${metrics['var_currency']:,.0f} | Expected Shortfall: ${metrics['es_currency']:,.0f}",
        f"Main modeled fragility driver: {driver}",
    ])
    allocation_lines = [
        f"{ticker}: {current:.1%} → {resilient:.1%} ({resilient-current:+.1%})"
        for ticker, current, resilient in zip(config.tickers, config.weights, result.resilient_weights)
    ]
    add_slide("Allocation", allocation_lines)
    add_slide("Performance and benchmark", [
        f"Portfolio return: {metrics['historical_return']:.1%} | {config.benchmark}: {metrics['benchmark_return']:.1%}",
        f"Portfolio volatility: {metrics['historical_volatility']:.1%} | {config.benchmark}: {metrics['benchmark_volatility']:.1%}",
        f"Sharpe: {metrics['sharpe']:.2f} | Beta: {metrics['beta']:.2f}",
        f"Maximum drawdown: {metrics['historical_max_drawdown']:.1%}",
    ])
    add_slide("Stress and fragility", [
        f"Worst modeled scenario: {worst['Scenario']}",
        f"Net shock: {worst['Net Portfolio Shock']:.1%} | Stressed value: ${worst['Stressed Value']:,.0f}",
        f"Crisis correlation: {metrics['normal_correlation']:.2f} → {metrics['crisis_correlation']:.2f}",
        f"Liquidity shortfall probability: {metrics['liquidity_shortfall_probability']:.1%}",
    ])
    add_slide("Methodology and limitations", [
        f"{config.method}; {config.simulations:,} paths; {config.simulation_days} trading days; {config.confidence:.0%} confidence",
        "Historical estimates, stylized stress shocks, and transparent liquidity proxies are used.",
        DISCLAIMER,
    ])
    output = BytesIO()
    prs.save(output)
    return output.getvalue()
